#!/usr/bin/env python3
"""Post-render PPTX geometry QA for Academic Guizang decks.

Checks the normalized PPTX itself (not only deck-spec) for:
- severe text-box overlaps;
- out-of-bounds shapes;
- non-cropped picture aspect-ratio distortion (especially logos/brand marks);
- collisions between small edge-brand images and text.

Intentional visual layering should be expressed with background/panel shapes rather than
by overlapping independent text boxes. Generated scientific figures are pictures and are
not inspected internally here; their scientific fidelity is validated upstream.
"""

from __future__ import annotations

import hashlib
import io
import math
import re
import sys
from collections import Counter
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as exc:  # pragma: no cover
    print(f"ERROR: python-pptx is required: {exc}", file=sys.stderr)
    raise SystemExit(2)

try:
    from PIL import Image
except Exception as exc:  # pragma: no cover
    print(f"ERROR: Pillow is required: {exc}", file=sys.stderr)
    raise SystemExit(2)

EMU = 914400.0


def ibox(shape):
    return (shape.left / EMU, shape.top / EMU, shape.width / EMU, shape.height / EMU)


def overlap_ratio(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    iw = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    ih = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    area = iw * ih
    if area <= 0:
        return 0.0
    return area / max(1e-9, min(aw * ah, bw * bh))


def text_of(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return (shape.text or "").strip()
    except Exception:
        return ""


def crop_is_zero(shape):
    vals = []
    for attr in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        try:
            vals.append(float(getattr(shape, attr)))
        except Exception:
            vals.append(0.0)
    return all(abs(v) < 1e-6 for v in vals)


def picture_info(shape):
    try:
        blob = shape.image.blob
        digest = hashlib.sha1(blob).hexdigest()
        im = Image.open(io.BytesIO(blob))
        natural_ratio = im.width / im.height
        x, y, w, h = ibox(shape)
        placed_ratio = w / h if h > 0 else float("inf")
        distortion = abs(math.log(max(placed_ratio, 1e-9) / max(natural_ratio, 1e-9)))
        return {
            "hash": digest,
            "natural_ratio": natural_ratio,
            "placed_ratio": placed_ratio,
            "distortion": distortion,
            "box": (x, y, w, h),
            "area": w * h,
            "crop_zero": crop_is_zero(shape),
        }
    except Exception:
        return None


ESTIMATE_RE = re.compile(r"\b\d+(?:\.\d+)?\s*\(\s*\d+(?:\.\d+)?\s*[–-]\s*\d+(?:\.\d+)?\s*\)")


def infer_picture_required(text, nontext_primitives):
    low = text.lower()
    estimate_count = len(ESTIMATE_RE.findall(text))
    if nontext_primitives >= 6 and estimate_count >= 4:
        return 'effect-plot'
    state_terms = sum(term in low for term in ['current use', 'non-use', 'no use', 'discontinuation', 'three-state', 'two-state'])
    if nontext_primitives >= 8 and state_terms >= 5:
        return 'treatment-state-diagram'
    trial_terms = sum(term in low for term in ['trial', 'eligible visit', 'time zero', 'clone', 'grace period', 'censor', 'ipcw', 'follow-up'])
    if nontext_primitives >= 15 and trial_terms >= 4:
        return 'study-design-diagram'
    pipeline_terms = sum(term in low for term in ['pipeline', 'workflow', 'cohort flow', 'stage', 'sequence'])
    if nontext_primitives >= 10 and pipeline_terms >= 2:
        return 'methods-pipeline'
    concept_terms = sum(term in low for term in ['estimand', 'causal', 'mechanism', 'conceptual', 'exposure state', 'treatment state'])
    if nontext_primitives >= 8 and concept_terms >= 2:
        return 'conceptual-schematic'
    return None


def is_large_figure_picture(info, sw, sh):
    x, y, w, h = info['box']
    # Small edge pictures are brand marks; a scientific figure should occupy meaningful body area.
    return (w * h >= 5.0 or (w >= sw * 0.35 and h >= sh * 0.22)) and not (w < 2.8 and h < 1.4)


def main():
    if len(sys.argv) != 2:
        print("Usage: python check-pptx-layout-integrity.py <deck.pptx>", file=sys.stderr)
        return 2
    pptx_path = Path(sys.argv[1]).resolve()
    if not pptx_path.exists():
        print(f"ERROR: not found: {pptx_path}", file=sys.stderr)
        return 2

    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width / EMU, prs.slide_height / EMU
    errors, warnings = [], []

    all_pictures = []
    per_slide = []
    for si, slide in enumerate(prs.slides, 1):
        texts = []
        pictures = []
        nontext_primitives = 0
        for shape in slide.shapes:
            x, y, w, h = ibox(shape)
            if x < -0.01 or y < -0.01 or x + w > sw + 0.01 or y + h > sh + 0.01:
                errors.append(f"Slide {si}: shape out of bounds: {shape.name} x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
            txt = text_of(shape)
            if txt:
                texts.append((shape, (x, y, w, h), txt))
            elif shape.shape_type not in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE}:
                nontext_primitives += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info = picture_info(shape)
                if info:
                    pictures.append((shape, info))
                    all_pictures.append((si, shape, info))
        per_slide.append((texts, pictures))

        combined_text = ' '.join(t for _, _, t in texts)
        inferred = infer_picture_required(combined_text, nontext_primitives)
        large_figure_pictures = [info for _, info in pictures if is_large_figure_picture(info, sw, sh)]
        if inferred and not large_figure_pictures:
            errors.append(
                f"Slide {si}: looks like whitelisted {inferred} ({nontext_primitives} non-text primitives) "
                "but contains no large scientific figure picture. Generate/render the illustration/plot as an image asset instead of native PowerPoint primitives."
            )

        # Severe text-text collisions. Background shapes are intentionally excluded.
        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                sa, ba, ta = texts[i]
                sb, bb, tb = texts[j]
                ratio = overlap_ratio(ba, bb)
                if ratio > 0.45:
                    errors.append(
                        f"Slide {si}: severe text overlap ratio={ratio:.2f}: {ta[:44]!r} <-> {tb[:44]!r}"
                    )
                elif ratio > 0.22:
                    warnings.append(
                        f"Slide {si}: possible text overlap ratio={ratio:.2f}: {ta[:44]!r} <-> {tb[:44]!r}"
                    )

    counts = Counter(info["hash"] for _, _, info in all_pictures)

    # Aspect-ratio distortion: a non-cropped picture should preserve its intrinsic ratio.
    for si, shape, info in all_pictures:
        if not info["crop_zero"]:
            continue
        x, y, w, h = info["box"]
        full_slide_visual = w >= sw * 0.93 and h >= sh * 0.90
        if info["distortion"] > 0.08:
            msg = (
                f"Slide {si}: picture aspect ratio distorted ({shape.name}): "
                f"placed={info['placed_ratio']:.3f}, intrinsic={info['natural_ratio']:.3f}. "
                f"Use contain/crop; never stretch logos."
            )
            if full_slide_visual:
                warnings.append(msg)
            else:
                errors.append(msg)

    # Small edge pictures are likely logos/brand marks: keep them clear of text.
    for si, (texts, pictures) in enumerate(per_slide, 1):
        for shape, info in pictures:
            x, y, w, h = info["box"]
            near_edge = x < 0.8 or y < 0.8 or x + w > sw - 0.8 or y + h > sh - 0.8
            likely_brand = info["area"] < 4.0 and near_edge and (counts[info["hash"]] >= 2 or info["natural_ratio"] > 1.5)
            if not likely_brand:
                continue
            for tshape, tbox, txt in texts:
                ratio = overlap_ratio(info["box"], tbox)
                if ratio > 0.05:
                    errors.append(f"Slide {si}: likely logo/brand image collides with text {txt[:44]!r} (ratio={ratio:.2f})")

    for w in warnings:
        print("WARNING:", w)
    if errors:
        for e in errors:
            print("ERROR:", e, file=sys.stderr)
        return 1
    print(f"PPTX layout integrity passed: {len(prs.slides)} slide(s).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
