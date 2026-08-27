#!/usr/bin/env python3
"""Report native PowerPoint text sizes for presentation readability QA.

This is a diagnostic tool, not a substitute for visual review. It inspects explicit
run/paragraph font sizes in text frames and tables. Figure-internal raster/SVG text
cannot be measured here and must follow figure-render-qa.md.
"""
from __future__ import annotations

import argparse
from collections import Counter
from pathlib import Path
from typing import Iterable

from pptx import Presentation


def iter_paragraphs(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for p in shape.text_frame.paragraphs:
                yield shape, p
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        yield shape, p


def paragraph_size_pt(p) -> float | None:
    sizes = [r.font.size.pt for r in p.runs if r.text.strip() and r.font.size is not None]
    if sizes:
        # Use the smallest explicit run size because a small embedded run can be the readability bottleneck.
        return min(sizes)
    if p.font.size is not None:
        return p.font.size.pt
    return None


def clean_text(p) -> str:
    return " ".join(r.text for r in p.runs).strip() or p.text.strip()


def main() -> int:
    ap = argparse.ArgumentParser(description="Report native PPTX text sizes for live-presentation readability QA.")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--target", type=float, default=14.0, help="Default audience-readable target in pt (default: 14).")
    ap.add_argument("--soft-min", type=float, default=12.0, help="Supporting-text soft minimum in pt (default: 12).")
    ap.add_argument("--hard-min", type=float, default=10.0, help="Fail threshold for non-footer native text (default: 10).")
    ap.add_argument("--footer-band", type=float, default=0.48, help="Bottom slide band in inches ignored for hard-min checks (default: 0.48).")
    ap.add_argument("--fail", action="store_true", help="Exit nonzero if any non-footer text is below hard-min.")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    slide_h = prs.slide_height / 914400.0
    counts = Counter()
    below_target = []
    below_soft = []
    below_hard = []
    unknown = 0

    for si, slide in enumerate(prs.slides, start=1):
        for shape, p in iter_paragraphs(slide):
            txt = clean_text(p)
            if not txt:
                continue
            size = paragraph_size_pt(p)
            if size is None:
                unknown += 1
                continue
            rounded = round(size * 2) / 2
            counts[rounded] += 1
            y = getattr(shape, "top", 0) / 914400.0
            h = getattr(shape, "height", 0) / 914400.0
            in_footer = (y + h) >= (slide_h - args.footer_band)
            rec = (si, size, txt[:120], in_footer)
            if size < args.target:
                below_target.append(rec)
            if size < args.soft_min:
                below_soft.append(rec)
            if size < args.hard_min and not in_footer:
                below_hard.append(rec)

    print(f"PPTX: {args.pptx}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Explicit-size text paragraphs: {sum(counts.values())}; unresolved/default-size paragraphs: {unknown}")
    print(f"Target audience-readable band begins at: {args.target:g} pt")
    print(f"Below target: {len(below_target)}; below soft minimum {args.soft_min:g} pt: {len(below_soft)}; below hard minimum {args.hard_min:g} pt outside footer: {len(below_hard)}")
    if counts:
        print("Size histogram (pt -> paragraph count):")
        print("  " + ", ".join(f"{k:g}:{v}" for k, v in sorted(counts.items())))

    if below_soft:
        print("\nExamples below soft minimum:")
        for si, size, txt, in_footer in below_soft[:40]:
            tag = "footer" if in_footer else "content"
            print(f"  slide {si:02d} | {size:g} pt | {tag} | {txt}")
        if len(below_soft) > 40:
            print(f"  ... {len(below_soft)-40} more")

    print("\nNote: raster/SVG figure-internal text is not measurable from PPTX native text. Review those figures using references/figure-render-qa.md and final on-slide equivalent size.")

    if args.fail and below_hard:
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
