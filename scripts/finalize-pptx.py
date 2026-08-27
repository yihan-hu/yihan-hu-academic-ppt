#!/usr/bin/env python3
"""Normalize and certify a PPTX before delivery.

The raw PptxGenJS package is never considered the final deliverable. This script:
1) validates the OOXML/ZIP package structurally;
2) opens it with python-pptx when available;
3) opens + re-saves it through LibreOffice/soffice (default, hard requirement);
4) validates the normalized package again and checks slide/notes/size invariants;
5) asks the office engine to render the normalized deck to PDF;
6) atomically writes the certified normalized PPTX to the requested output path.

Use --structural-only only when no office engine exists and the user explicitly accepts
that cross-application openability cannot be certified.
"""

from __future__ import annotations

import argparse
import os
import posixpath
import re
import shutil
import subprocess
import sys
import tempfile
import urllib.parse
import zipfile
from dataclasses import dataclass
from pathlib import Path
import xml.etree.ElementTree as ET

RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REQUIRED_PARTS = {
    "[Content_Types].xml",
    "_rels/.rels",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
}


class ValidationError(RuntimeError):
    pass


@dataclass
class PackageMetrics:
    slides: int
    notes: int
    width_emu: int | None
    height_emu: int | None


def _relationship_base_dir(rels_path: str) -> str:
    if rels_path == "_rels/.rels":
        return ""
    marker = "/_rels/"
    if marker not in rels_path or not rels_path.endswith(".rels"):
        raise ValidationError(f"unexpected relationships path: {rels_path}")
    prefix, rel_name = rels_path.split(marker, 1)
    source_name = rel_name[:-5]
    if not source_name:
        raise ValidationError(f"invalid relationships filename: {rels_path}")
    return prefix


def _normalize_rel_target(base_dir: str, target: str) -> str | None:
    target = target.split("#", 1)[0]
    target = urllib.parse.unquote(target)
    if not target:
        return None
    parsed = urllib.parse.urlparse(target)
    if parsed.scheme and parsed.scheme not in {"file"}:
        return None
    if target.startswith("/"):
        normalized = posixpath.normpath(target.lstrip("/"))
    else:
        normalized = posixpath.normpath(posixpath.join(base_dir, target))
    if normalized.startswith("../") or normalized == "..":
        raise ValidationError(f"relationship target escapes package root: {target}")
    return normalized


def _parse_slide_size(raw: bytes) -> tuple[int | None, int | None]:
    root = ET.fromstring(raw)
    sld_sz = root.find(f"{{{P_NS}}}sldSz")
    if sld_sz is None:
        return None, None
    try:
        return int(sld_sz.attrib.get("cx", "")), int(sld_sz.attrib.get("cy", ""))
    except ValueError:
        return None, None


def validate_package(path: Path, *, label: str) -> PackageMetrics:
    if not path.exists() or path.stat().st_size == 0:
        raise ValidationError(f"{label}: PPTX does not exist or is empty: {path}")

    try:
        with zipfile.ZipFile(path) as zf:
            names = zf.namelist()
            if len(names) != len(set(names)):
                dupes = sorted({n for n in names if names.count(n) > 1})
                raise ValidationError(f"{label}: duplicate ZIP members: {dupes[:8]}")

            bad_member = zf.testzip()
            if bad_member:
                raise ValidationError(f"{label}: CRC failure in ZIP member: {bad_member}")

            missing = sorted(REQUIRED_PARTS - set(names))
            if missing:
                raise ValidationError(f"{label}: missing required OOXML part(s): {', '.join(missing)}")

            xml_parts = [n for n in names if n.endswith((".xml", ".rels"))]
            for name in xml_parts:
                try:
                    ET.fromstring(zf.read(name))
                except ET.ParseError as exc:
                    raise ValidationError(f"{label}: malformed XML in {name}: {exc}") from exc

            name_set = set(names)
            for rels_path in [n for n in names if n.endswith(".rels")]:
                base_dir = _relationship_base_dir(rels_path)
                root = ET.fromstring(zf.read(rels_path))
                for rel in root.findall(f"{{{RELS_NS}}}Relationship"):
                    if rel.attrib.get("TargetMode") == "External":
                        continue
                    target = rel.attrib.get("Target", "")
                    normalized = _normalize_rel_target(base_dir, target)
                    if normalized and normalized not in name_set:
                        raise ValidationError(
                            f"{label}: broken relationship in {rels_path}: {target} -> {normalized} not found"
                        )

            slide_parts = [n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)]
            notes_parts = [n for n in names if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", n)]
            width, height = _parse_slide_size(zf.read("ppt/presentation.xml"))
            return PackageMetrics(len(slide_parts), len(notes_parts), width, height)
    except zipfile.BadZipFile as exc:
        raise ValidationError(f"{label}: not a valid ZIP/PPTX package: {exc}") from exc


def validate_with_python_pptx(path: Path, *, expected_slides: int | None = None) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        print("WARNING: python-pptx is unavailable; skipping independent parser check.", file=sys.stderr)
        return
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ValidationError(f"python-pptx could not open {path.name}: {exc}") from exc
    if expected_slides is not None and len(prs.slides) != expected_slides:
        raise ValidationError(
            f"python-pptx sees {len(prs.slides)} slides, expected {expected_slides}"
        )


def find_office_engine() -> str | None:
    for name in ("libreoffice", "soffice"):
        exe = shutil.which(name)
        if exe:
            return exe
    return None


def run_office_convert(exe: str, source: Path, outdir: Path, fmt: str, profile: Path) -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    profile.mkdir(parents=True, exist_ok=True)
    cmd = [
        exe,
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        f"-env:UserInstallation={profile.resolve().as_uri()}",
        "--convert-to",
        fmt,
        "--outdir",
        str(outdir),
        str(source),
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, timeout=180)
    log = (proc.stdout or "").strip()
    if proc.returncode != 0:
        raise ValidationError(f"office conversion failed ({proc.returncode}): {log}")

    ext = fmt.split(":", 1)[0].strip().lower()
    expected = outdir / f"{source.stem}.{ext}"
    if expected.exists() and expected.stat().st_size > 0:
        return expected

    candidates = sorted(outdir.glob(f"*.{ext}"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not candidates:
        raise ValidationError(f"office conversion reported success but produced no .{ext}: {log}")
    return candidates[0]


def compare_metrics(raw: PackageMetrics, normalized: PackageMetrics) -> None:
    if normalized.slides != raw.slides:
        raise ValidationError(
            f"normalization changed slide count: raw={raw.slides}, normalized={normalized.slides}"
        )
    if raw.notes > 0 and normalized.notes < raw.notes:
        raise ValidationError(
            f"normalization dropped notes slides: raw={raw.notes}, normalized={normalized.notes}"
        )
    if raw.width_emu and raw.height_emu and normalized.width_emu and normalized.height_emu:
        if (raw.width_emu, raw.height_emu) != (normalized.width_emu, normalized.height_emu):
            raise ValidationError(
                "normalization changed slide size: "
                f"raw={raw.width_emu}x{raw.height_emu}, "
                f"normalized={normalized.width_emu}x{normalized.height_emu}"
            )


def atomic_copy(source: Path, destination: Path) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp_name = tempfile.mkstemp(prefix=destination.name + ".", suffix=".tmp", dir=str(destination.parent))
    os.close(fd)
    tmp = Path(tmp_name)
    try:
        shutil.copy2(source, tmp)
        os.replace(tmp, destination)
    finally:
        if tmp.exists():
            tmp.unlink()


def run_post_qa(final_pptx: Path, brand_profile: str | None) -> None:
    script_dir = Path(__file__).resolve().parent
    layout = script_dir / "check-pptx-layout-integrity.py"
    if layout.exists():
        proc = subprocess.run([sys.executable, str(layout), str(final_pptx)], stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
        if proc.returncode != 0:
            raise ValidationError("post-normalization layout/figure integrity failed: " + (proc.stdout or '').strip().replace('\n', ' | '))
        if (proc.stdout or '').strip():
            print((proc.stdout or '').strip())
    if brand_profile in {"ki-editorial", "ki-swiss", "ki"}:
        pal = script_dir / "check-ki-pptx-palette.py"
        if pal.exists():
            profile = "ki-swiss" if brand_profile in {"ki-swiss", "ki"} else "ki-editorial"
            proc = subprocess.run([sys.executable, str(pal), str(final_pptx), "--profile", profile], stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
            if proc.returncode != 0:
                raise ValidationError("post-normalization KI palette audit failed: " + (proc.stdout or '').strip().replace('\n', ' | '))
            if (proc.stdout or '').strip():
                print((proc.stdout or '').strip())


def main() -> int:
    parser = argparse.ArgumentParser(description="Normalize and certify a PPTX before delivery.")
    parser.add_argument("input_pptx", type=Path, help="raw PPTX generated by PptxGenJS")
    parser.add_argument("output_pptx", type=Path, help="final certified PPTX to deliver")
    parser.add_argument(
        "--structural-only",
        action="store_true",
        help="skip office-engine open/resave/render; does NOT certify real application openability",
    )
    parser.add_argument(
        "--brand-profile",
        default=None,
        help="optional bundled brand profile id; KI profiles trigger final palette audit",
    )
    args = parser.parse_args()

    source = args.input_pptx.resolve()
    destination = args.output_pptx.resolve()

    try:
        raw_metrics = validate_package(source, label="raw")
        validate_with_python_pptx(source, expected_slides=raw_metrics.slides)
        print(
            f"Raw structural check passed: slides={raw_metrics.slides}, notes={raw_metrics.notes}, "
            f"size={source.stat().st_size} bytes"
        )

        if args.structural_only:
            atomic_copy(source, destination)
            run_post_qa(destination, args.brand_profile)
            print("WARNING: structural-only mode used; application openability was NOT certified.", file=sys.stderr)
            print(f"Structural-only output written: {destination}")
            return 0

        office = find_office_engine()
        if not office:
            raise ValidationError(
                "LibreOffice/soffice not found. Refusing to certify a final PPTX without a real office-engine open/resave test. "
                "Use --structural-only only if the user explicitly accepts the compatibility risk."
            )

        with tempfile.TemporaryDirectory(prefix="academic-ppt-finalize-") as td:
            td_path = Path(td)
            normalized_dir = td_path / "normalized"
            render_dir = td_path / "render"
            normalized = run_office_convert(
                office,
                source,
                normalized_dir,
                "pptx:Impress Office Open XML",
                td_path / "profile-normalize",
            )
            normalized_metrics = validate_package(normalized, label="normalized")
            compare_metrics(raw_metrics, normalized_metrics)
            validate_with_python_pptx(normalized, expected_slides=raw_metrics.slides)

            pdf = run_office_convert(
                office,
                normalized,
                render_dir,
                "pdf",
                td_path / "profile-render",
            )
            if pdf.stat().st_size < 1024:
                raise ValidationError(f"office-engine PDF render is suspiciously small: {pdf.stat().st_size} bytes")

            atomic_copy(normalized, destination)

        final_metrics = validate_package(destination, label="final")
        compare_metrics(raw_metrics, final_metrics)
        validate_with_python_pptx(destination, expected_slides=raw_metrics.slides)
        run_post_qa(destination, args.brand_profile)
        print(
            "PPTX finalization passed: raw OOXML -> office open/resave -> normalized OOXML -> office PDF render -> final OOXML -> layout/figure QA -> brand palette QA when applicable."
        )
        print(f"Final certified PPTX: {destination}")
        return 0
    except (ValidationError, subprocess.TimeoutExpired, OSError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
